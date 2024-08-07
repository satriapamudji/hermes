import os
import docx
import PyPDF2
import logging
import pdfplumber
from pptx import Presentation

from .helpers import download_file, get_file_extension, get_file_name

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        text.append(' '.join(slide_text))

    full_text = '\n\n'.join(text)
    logger.info(f"Extracted {len(full_text)} characters from PPTX with {len(prs.slides)} slides")
    return full_text

def extract_text_from_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ' '.join(page.extract_text() for page in reader.pages)
        logger.info(f"Extracted {len(text)} characters from PDF with {len(reader.pages)} pages using PyPDF2")
        return text
    except Exception as e:
        logger.warning(f"PyPDF2 failed to extract text from PDF: {str(e)}")
        logger.info("Trying to extract text using pdfplumber...")
        try:
            with pdfplumber.open(file_path) as pdf:
                text = ' '.join(page.extract_text() for page in pdf.pages)
            logger.info(f"Extracted {len(text)} characters from PDF with {len(pdf.pages)} pages using pdfplumber")
            return text
        except Exception as e:
            logger.error(f"pdfplumber also failed to extract text from PDF: {str(e)}", exc_info=True)
            return ""

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = ' '.join(paragraph.text for paragraph in doc.paragraphs)
    logger.info(f"Extracted {len(text)} characters from DOCX with {len(doc.paragraphs)} paragraphs")
    return text

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    logger.info(f"Extracted {len(text)} characters from TXT file")
    return text

async def extract_text_from_file(client, message, file, doc_number):
    user_id = message.sender_id
    timestamp = message.date
    logger.info(f"{timestamp} - Starting text extraction for file {doc_number} from user {user_id}")

    file_path = await download_file(client, message, file)
    extension = get_file_extension(file).lower()
    file_name = get_file_name(file)
    
    logger.info(f"{timestamp} - File {doc_number} downloaded for user {user_id}. File name: {file_name}, Extension: {extension}")

    try:
        if extension == '.pdf':
            logger.info(f"{timestamp} - Extracting text from PDF file {doc_number} for user {user_id}")
            text = extract_text_from_pdf(file_path)
        elif extension == '.docx':
            logger.info(f"{timestamp} - Extracting text from DOCX file {doc_number} for user {user_id}")
            text = extract_text_from_docx(file_path)
        elif extension == '.txt':
            logger.info(f"{timestamp} - Extracting text from TXT file {doc_number} for user {user_id}")
            text = extract_text_from_txt(file_path)
        elif extension == '.pptx':
            logger.info(f"{timestamp} - Extracting text from PPTX file {doc_number} for user {user_id}")
            text = extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")

        logger.info(f"{timestamp} - Text extraction completed for file {doc_number} from user {user_id}")
        logger.info(f"{timestamp} - Extracted {len(text)} characters from file {doc_number} for user {user_id}")
        logger.debug(f"{timestamp} - First 500 characters of extracted text from file {doc_number} for user {user_id}: {text[:500]}")

    except Exception as e:
        logger.error(f"{timestamp} - Error extracting text from file {doc_number} for user {user_id}: {str(e)}", exc_info=True)
        raise

    finally:
        # Clean up the temporary file
        os.remove(file_path)
        logger.info(f"{timestamp} - Temporary file for document {doc_number} removed for user {user_id}")
    
    return {
        'doc_number': doc_number,
        'file_name': file_name,
        'text': text
    }