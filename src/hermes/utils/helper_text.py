import os
import docx
import json
import PyPDF2
import tempfile
import pdfplumber
from pptx import Presentation
from telethon.tl.types import DocumentAttributeFilename
from logging_config import logger

def read_json_file(file_name):
    file_path = os.path.join('data', file_name)
    logger.info(f"Reading JSON file from {file_path}")
    with open(file_path, 'r', encoding='utf-8') as file:
        return json.load(file)

def read_text_file(file_name):
    file_path = os.path.join('data', file_name)
    logger.info(f"Reading text file from {file_path}")
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

async def download_text_file(client, message, file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=get_file_extension(file)) as temp_file:
        await client.download_media(file, temp_file.name)
        logger.info(f"Downloaded file to {temp_file.name}")
        return temp_file.name

def get_file_extension(file):
    for attr in file.document.attributes:
        if isinstance(attr, DocumentAttributeFilename):
            return os.path.splitext(attr.file_name)[1]
    return ''

def get_file_name(file):
    for attr in file.document.attributes:
        if isinstance(attr, DocumentAttributeFilename):
            return attr.file_name
    return 'Unnamed Document'

def extract_text_from_pptx(file_path):
    logger.info(f"Extracting text from PPTX file at {file_path}")
    prs = Presentation(file_path)
    text = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text)
        text.append(' '.join(slide_text))

    return '\n\n'.join(text)

def extract_text_from_pdf(file_path):
    logger.info(f"Extracting text from PDF file at {file_path}")
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            return ' '.join(page.extract_text() for page in reader.pages)
    except Exception:
        logger.warning("PyPDF2 failed to extract text, attempting with pdfplumber", exc_info=True)
        try:
            with pdfplumber.open(file_path) as pdf:
                return ' '.join(page.extract_text() for page in pdf.pages)
        except Exception as e:
            logger.error(f"Failed to extract text from PDF using both PyPDF2 and pdfplumber: {str(e)}", exc_info=True)
            raise RuntimeError(f"Failed to extract text from PDF: {str(e)}")

def extract_text_from_docx(file_path):
    logger.info(f"Extracting text from DOCX file at {file_path}")
    doc = docx.Document(file_path)
    return ' '.join(paragraph.text for paragraph in doc.paragraphs)

def extract_text_from_txt(file_path):
    logger.info(f"Extracting text from TXT file at {file_path}")
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

async def extract_text_from_file(client, message, file, doc_number):
    file_path = await download_text_file(client, message, file)
    extension = get_file_extension(file).lower()
    file_name = get_file_name(file)
    logger.info(f"Processing file {file_name} with extension {extension}")

    try:
        if extension == '.pdf':
            text = extract_text_from_pdf(file_path)
        elif extension == '.docx':
            text = extract_text_from_docx(file_path)
        elif extension == '.txt':
            text = extract_text_from_txt(file_path)
        elif extension == '.pptx':
            text = extract_text_from_pptx(file_path)
        else:
            logger.error(f"Unsupported file type: {extension}")
            raise ValueError(f"Unsupported file type: {extension}")

        logger.info(f"Text extraction completed for file {doc_number} ({file_name})")
    except Exception as e:
        logger.error(f"Error extracting text from file {doc_number} ({file_name}): {str(e)}", exc_info=True)
        raise
    finally:
        os.remove(file_path)
        logger.info(f"Temporary file {file_path} removed")
    
    return {
        'doc_number': doc_number,
        'file_name': file_name,
        'text': text
    }
